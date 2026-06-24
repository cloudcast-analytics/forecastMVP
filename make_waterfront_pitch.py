"""
CloudCast x Waterfront Genk — Sales pitch voor Kris
Dinsdag 10u, ter plaatse
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Kleuren ────────────────────────────────────────────────────────────────
BLUE        = RGBColor(0x1a, 0x44, 0xe8)
BLUE_DARK   = RGBColor(0x0f, 0x2a, 0x9e)
BLUE_MID    = RGBColor(0x2d, 0x55, 0xf0)
BLUE_LIGHT  = RGBColor(0xe8, 0xee, 0xff)
BLUE_PALE   = RGBColor(0xf0, 0xf4, 0xff)
TEXT        = RGBColor(0x1a, 0x1f, 0x36)
MUTED       = RGBColor(0x6b, 0x72, 0x80)
WHITE       = RGBColor(0xff, 0xff, 0xff)
GREEN       = RGBColor(0x16, 0xa3, 0x4a)
GREEN_LIGHT = RGBColor(0xdc, 0xfc, 0xe7)
AMBER       = RGBColor(0xd9, 0x77, 0x06)
RED         = RGBColor(0xdc, 0x26, 0x26)
WATER_BLUE  = RGBColor(0x03, 0x69, 0xa1)   # Waterfront accent

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Hulpfuncties ───────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt or 1)
    return s


def txt(slide, text, x, y, w, h,
        size=16, bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, font="Calibri", wrap=True):
    color = color or TEXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def txt2(slide, lines, x, y, w, h, size=14, font="Calibri"):
    """lines = list of (text, bold, color)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def circ(slide, x, y, d, fill):
    s = slide.shapes.add_shape(9, x, y, d, d)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def dot(slide, x, y, color=BLUE):
    circ(slide, x, y, Inches(0.1), color)


def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.35), fill=BLUE)
    rect(slide, 0, Inches(1.35), W, Pt(4), fill=BLUE_MID)
    txt(slide, title, Inches(0.65), Inches(0.18), Inches(11.5), Inches(0.7),
        size=30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.65), Inches(0.88), Inches(11.5), Inches(0.38),
            size=14, color=RGBColor(0xc7, 0xd4, 0xfd))


def card(slide, x, y, w, h, fill=WHITE,
         border=None, top_bar_color=None, top_bar_h=Inches(0.42)):
    c = rect(slide, x, y, w, h, fill=fill,
             line_color=border or RGBColor(0xe2, 0xe8, 0xf0), line_pt=1)
    if top_bar_color:
        rect(slide, x, y, w, top_bar_h, fill=top_bar_color)
    return c


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# Boodschap: Dit gaat over jou, Kris.
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=RGBColor(0xf0, 0xf4, 0xff))

# Linker blauwe balk (1/3 van het scherm)
rect(sl, 0, 0, Inches(5.2), H, fill=BLUE)

# Waterfront decoratieve balk onderaan links
rect(sl, 0, Inches(6.5), Inches(5.2), Inches(1.0),
     fill=RGBColor(0x0f, 0x2a, 0x9e))

# Logo tekst links
txt(sl, "CloudCast", Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.8),
    size=38, bold=True, color=WHITE)
txt(sl, "Analytics", Inches(0.5), Inches(1.95), Inches(4.2), Inches(0.6),
    size=28, color=RGBColor(0xc7, 0xd4, 0xfd))

rect(sl, Inches(0.5), Inches(2.75), Inches(2.8), Pt(3), fill=WHITE)

txt(sl, "Revenue Forecasting\nvoor Horeca & Leisure",
    Inches(0.5), Inches(2.95), Inches(4.2), Inches(1.0),
    size=14, color=RGBColor(0xc7, 0xd4, 0xfd), italic=True)

txt(sl, "cloudcastanalytics.com",
    Inches(0.5), Inches(6.6), Inches(4.2), Inches(0.4),
    size=12, color=RGBColor(0x93, 0xc5, 0xfd))

# Rechts: de hook
txt(sl, "Kris,",
    Inches(5.6), Inches(1.0), Inches(7.2), Inches(0.75),
    size=40, bold=True, color=TEXT)

txt(sl, "stel je voor dat je\nmorgenvroeg al weet\nhoe druk het\nvolgende zaterdag\nwordt.",
    Inches(5.6), Inches(1.8), Inches(7.0), Inches(3.5),
    size=30, color=TEXT)

txt(sl, "Dat is precies wat CloudCast doet.",
    Inches(5.6), Inches(5.4), Inches(7.0), Inches(0.5),
    size=16, italic=True, color=MUTED)

# Datum + info
txt(sl, "Dinsdag 24 juni 2026  |  Waterfront Genk",
    Inches(5.6), Inches(6.7), Inches(7.0), Inches(0.4),
    size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — HET PROBLEEM
# Boodschap: Jullie verliezen geld elke week. Niet door slechte service — door slechte planning.
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Jullie verliezen geld. Niet door slechte service.",
           "Maar door één ding: plannen op gevoel.")

# Grote pijnpunt-kaarten
pijnen = [
    ("Te veel personeel\nop een rustige dag",
     "Elke te veel ingeplande medewerker kost €110–€150 per dag.\n"
     "Op een stille maandag in oktober zijn dat snel 3 FTE te veel.",
     "€ 330–450\nverlies per dag",
     RED),
    ("Te weinig personeel\nop een drukke dag",
     "Zonnige zaterdag in juli? Iedereen staat in de rij.\n"
     "Klanten haken af. Reviews dalen. Omzet blijft liggen.",
     "Omzet\nmisgelopen",
     AMBER),
    ("Niemand weet het\nop voorhand",
     "Roosters worden opgesteld op buikgevoel of vorig jaar.\n"
     "Maar vorig jaar was het anders. Het weer was anders. Het was een andere week.",
     "60–70%\nnauwkeurig",
     MUTED),
]

for i, (title, body, stat, col) in enumerate(pijnen):
    cx = Inches(0.45) + i * Inches(4.25)
    cy = Inches(1.65)
    cw = Inches(4.0)
    ch = Inches(5.3)

    card(sl, cx, cy, cw, ch,
         fill=RGBColor(0xf8, 0xf9, 0xff),
         border=RGBColor(0xe2, 0xe8, 0xf0))

    # Gekleurde bovenbalk
    rect(sl, cx, cy, cw, Inches(0.06), fill=col)

    txt(sl, title, cx + Inches(0.25), cy + Inches(0.2), cw - Inches(0.5), Inches(0.8),
        size=17, bold=True, color=TEXT)

    txt(sl, body, cx + Inches(0.25), cy + Inches(1.1), cw - Inches(0.5), Inches(1.8),
        size=13, color=MUTED)

    # Stat block onderaan de kaart
    sbox = rect(sl, cx + Inches(0.25), cy + Inches(3.5), cw - Inches(0.5), Inches(1.4),
                fill=col if col != MUTED else RGBColor(0xf1, 0xf5, 0xf9),
                line_color=None)
    txt(sl, stat,
        cx + Inches(0.3), cy + Inches(3.55), cw - Inches(0.6), Inches(1.3),
        size=22, bold=True,
        color=WHITE if col != MUTED else TEXT,
        align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CLOUDCAST IN ACTIE (demo-preview)
# Boodschap: Dit is wat jij elke ochtend zou zien.
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Dit is wat jij elke ochtend ziet.",
           "CloudCast — forecast voor Waterfront Genk, komende 7 dagen")

# "Browser-venster" mockup
rect(sl, Inches(0.45), Inches(1.55), Inches(12.4), Inches(5.5),
     fill=RGBColor(0xf8, 0xf9, 0xff),
     line_color=RGBColor(0xe2, 0xe8, 0xf0), line_pt=1)

# Topbalk van browser (nep)
rect(sl, Inches(0.45), Inches(1.55), Inches(12.4), Inches(0.45),
     fill=RGBColor(0xf1, 0xf5, 0xf9))
txt(sl, "  cloudcastanalytics.com/forecast  —  Waterfront Genk",
    Inches(0.55), Inches(1.6), Inches(8), Inches(0.32),
    size=10, color=MUTED)

# Forecast tabel header
cols = ["Dag", "Datum", "Weer", "Bezoekers", "Omzet", "FTE nodig", "Drukteniveau"]
col_widths = [Inches(1.0), Inches(1.3), Inches(1.0), Inches(1.4), Inches(1.6), Inches(1.5), Inches(1.8)]
col_x = Inches(0.55)
row_y = Inches(2.15)

rect(sl, Inches(0.55), row_y, Inches(12.2), Inches(0.38),
     fill=BLUE)

cx = col_x
for i, (col_name, cw) in enumerate(zip(cols, col_widths)):
    txt(sl, col_name, cx + Inches(0.05), row_y + Inches(0.06), cw - Inches(0.1), Inches(0.28),
        size=11, bold=True, color=WHITE)
    cx += cw

# Forecast data
forecast_rows = [
    ("Wo", "25 jun", "⛅ 21°C", "340",   "€ 11.900",  "9 FTE",  "Normaal",   MUTED,  RGBColor(0xf1, 0xf5, 0xf9)),
    ("Do", "26 jun", "☀️ 26°C", "510",   "€ 17.850",  "13 FTE", "Druk",      AMBER,  WHITE),
    ("Vr", "27 jun", "☀️ 28°C", "680",   "€ 23.800",  "17 FTE", "Druk",      AMBER,  WHITE),
    ("Za", "28 jun", "☀️ 29°C", "890",   "€ 31.150",  "22 FTE", "Heel druk", RED,    RGBColor(0xff, 0xf5, 0xf5)),
    ("Zo", "29 jun", "☀️ 28°C", "840",   "€ 29.400",  "21 FTE", "Heel druk", RED,    RGBColor(0xff, 0xf5, 0xf5)),
    ("Ma", "30 jun", "⛅ 22°C", "280",   "€ 9.800",   "7 FTE",  "Normaal",   MUTED,  RGBColor(0xf1, 0xf5, 0xf9)),
    ("Di", "1 jul",  "🌧️ 17°C", "190",   "€ 6.650",   "5 FTE",  "Rustig",    GREEN,  RGBColor(0xf0, 0xfd, 0xf4)),
]

for ri, (dag, datum, weer, bez, omzet, fte, niveau, niveau_col, row_fill) in enumerate(forecast_rows):
    ry = row_y + Inches(0.38) + ri * Inches(0.58)
    rect(sl, Inches(0.55), ry, Inches(12.2), Inches(0.56), fill=row_fill)
    if ri > 0:
        rect(sl, Inches(0.55), ry, Inches(12.2), Pt(1),
             fill=RGBColor(0xe2, 0xe8, 0xf0))

    row_vals = [dag, datum, weer, bez, omzet, fte, niveau]
    cx = col_x
    for vi, (val, cw) in enumerate(zip(row_vals, col_widths)):
        is_niveau = vi == 6
        val_color = TEXT
        bold = False
        if vi == 4:  # omzet
            val_color = BLUE; bold = True
        if is_niveau:
            val_color = niveau_col; bold = True
        txt(sl, val, cx + Inches(0.07), ry + Inches(0.12), cw - Inches(0.1), Inches(0.36),
            size=12, bold=bold, color=val_color)
        cx += cw


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CONCREET VOORBEELD: ZATERDAG 28 JUNI
# Boodschap: Dit is wat je donderdag al weet over aanstaande zaterdag.
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Zaterdag 28 juni — je weet het al op donderdag.",
           "Concreet voorbeeld: wat CloudCast je zou tonen voor aanstaande zaterdag")

# Grote omzet-kaart links
rect(sl, Inches(0.45), Inches(1.65), Inches(4.8), Inches(5.3), fill=BLUE)
txt(sl, "Verwachte omzet",
    Inches(0.7), Inches(1.95), Inches(4.3), Inches(0.45),
    size=14, color=RGBColor(0xc7, 0xd4, 0xfd))
txt(sl, "€ 31.150",
    Inches(0.7), Inches(2.45), Inches(4.3), Inches(1.0),
    size=42, bold=True, color=WHITE)
txt(sl, "Verwachte bezoekers",
    Inches(0.7), Inches(3.65), Inches(4.3), Inches(0.4),
    size=13, color=RGBColor(0xc7, 0xd4, 0xfd))
txt(sl, "890 personen",
    Inches(0.7), Inches(4.0), Inches(4.3), Inches(0.55),
    size=26, bold=True, color=WHITE)
txt(sl, "Aanbevolen personeel",
    Inches(0.7), Inches(4.75), Inches(4.3), Inches(0.4),
    size=13, color=RGBColor(0xc7, 0xd4, 0xfd))
txt(sl, "22 FTE",
    Inches(0.7), Inches(5.1), Inches(4.3), Inches(0.5),
    size=26, bold=True, color=WHITE)
txt(sl, "Terras: 10  |  Bar: 8  |  Waterski: 4",
    Inches(0.7), Inches(5.65), Inches(4.3), Inches(0.35),
    size=12, color=RGBColor(0x93, 0xc5, 0xfd))

# Rechts: factoren
txt(sl, "Waarom zo druk?",
    Inches(5.55), Inches(1.65), Inches(7.3), Inches(0.45),
    size=16, bold=True, color=TEXT)

factoren = [
    (GREEN,  "Zonnig weekend — 29°C",                        "Historisch: +42% bezoekers vs. bewolkt"),
    (GREEN,  "Schoolvakantie Limburg start",                  "Vorig jaar +28% op eerste vakantiezaterdag"),
    (BLUE,   "Geen groot evenement in Genk",                  "Geen afleiding van doelgroep"),
    (AMBER,  "Vrijdag ook druk — mensen blijven in de buurt", "Cross-dag effect: +12%"),
]

for i, (col, factor, toelichting) in enumerate(factoren):
    fy = Inches(2.25) + i * Inches(1.1)
    rect(sl, Inches(5.55), fy, Inches(7.3), Inches(0.92),
         fill=RGBColor(0xf8, 0xf9, 0xff),
         line_color=RGBColor(0xe2, 0xe8, 0xf0))
    rect(sl, Inches(5.55), fy, Inches(0.08), Inches(0.92), fill=col)
    txt(sl, factor,
        Inches(5.75), fy + Inches(0.08), Inches(7.0), Inches(0.36),
        size=13, bold=True, color=TEXT)
    txt(sl, toelichting,
        Inches(5.75), fy + Inches(0.46), Inches(7.0), Inches(0.35),
        size=12, color=MUTED)

txt(sl, "Betrouwbaarheidsinterval: € 28.000 – € 34.300",
    Inches(5.55), Inches(6.6), Inches(7.3), Inches(0.38),
    size=12, color=MUTED, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROI VOOR WATERFRONT GENK
# Boodschap: Dit is hoeveel het jullie oplevert. Specifiek.
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Wat levert het Waterfront Genk concreet op?",
           "Op basis van jullie type locatie en seizoenspatroon — geen theorie")

# Berekening uitgelegd links
txt(sl, "Hoe we dit berekenen:",
    Inches(0.5), Inches(1.65), Inches(6), Inches(0.4),
    size=14, bold=True, color=TEXT)

stappen_roi = [
    ("Gemiddeld 2 FTE te veel op drukke dagen",
     "Op 60% van de weekenddagen plant een zaak als Waterfront\n"
     "op gevoel 2 medewerkers te veel in."),
    ("Loonkost: €130/medewerker/dag",
     "Inclusief sociale lasten — realistisch voor horecapersoneel\n"
     "in Limburg."),
    ("52 weekends × 60% × 2 FTE × €130",
     "= €8.112 besparing op weekends alleen."),
    ("Plus weekdagwinst (40% van 220 dagen)",
     "Minder overbezetting op rustige weekdagen:\n"
     "+ €4.576 per jaar."),
]

for i, (title, body) in enumerate(stappen_roi):
    sy = Inches(2.15) + i * Inches(1.15)
    circ(sl, Inches(0.5), sy + Inches(0.06), Inches(0.38), BLUE)
    txt(sl, str(i + 1),
        Inches(0.5), sy + Inches(0.06), Inches(0.38), Inches(0.38),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title,
        Inches(1.05), sy, Inches(5.4), Inches(0.35),
        size=13, bold=True, color=TEXT)
    txt(sl, body,
        Inches(1.05), sy + Inches(0.36), Inches(5.4), Inches(0.72),
        size=12, color=MUTED)

# Groot resultaat rechts
rect(sl, Inches(7.0), Inches(1.65), Inches(5.85), Inches(5.3), fill=BLUE)

txt(sl, "Totale jaarlijkse besparing",
    Inches(7.25), Inches(1.95), Inches(5.35), Inches(0.45),
    size=14, color=RGBColor(0xc7, 0xd4, 0xfd))
txt(sl, "€ 12.688",
    Inches(7.25), Inches(2.45), Inches(5.35), Inches(1.1),
    size=48, bold=True, color=WHITE)
txt(sl, "per jaar",
    Inches(7.25), Inches(3.5), Inches(5.35), Inches(0.4),
    size=16, color=RGBColor(0xc7, 0xd4, 0xfd))

rect(sl, Inches(7.25), Inches(4.05), Inches(5.1), Pt(2),
     fill=RGBColor(0x93, 0xc5, 0xfd))

txt(sl, "CloudCast kost € 249/maand\n(Groei-plan, 1 locatie)",
    Inches(7.25), Inches(4.2), Inches(5.1), Inches(0.65),
    size=13, color=RGBColor(0xc7, 0xd4, 0xfd))

txt(sl, "Terugverdientijd:",
    Inches(7.25), Inches(5.0), Inches(5.1), Inches(0.35),
    size=14, color=RGBColor(0xc7, 0xd4, 0xfd))
txt(sl, "3 tot 4 weken",
    Inches(7.25), Inches(5.35), Inches(5.1), Inches(0.6),
    size=28, bold=True, color=RGBColor(0x86, 0xef, 0xac))

txt(sl, "Netto jaarwinst: > € 9.700",
    Inches(7.25), Inches(6.1), Inches(5.1), Inches(0.45),
    size=14, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOE WE STARTEN
# Boodschap: Dit is simpel. Geen IT-project. Klaar in 1 week.
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Van handdruk tot eerste forecast: 1 week.",
           "Geen IT-afdeling nodig. Geen technische kennis. Wij doen het zware werk.")

stappen = [
    ("Week 1\ndag 1–2",
     "Koppeling of upload",
     "Je stuurt ons een Excel-export van je kassasysteem "
     "(of we koppelen rechtstreeks). Wij verwerken alles. "
     "Jij hoeft niks technisch te doen."),
    ("Week 1\ndag 3–5",
     "Model leert",
     "CloudCast analyseert je historische data: "
     "seizoenen, weekenden, zonnige zomerdagen, lokale events. "
     "Na 3 dagen zijn de eerste patronen zichtbaar."),
    ("Vanaf\nweek 2",
     "Dagelijkse forecast",
     "Elke ochtend een nieuwe 14-daagse forecast. "
     "Omzet, bezoekers, aanbevolen personeel per dag. "
     "In je dashboard. Klaar om op te handelen."),
]

step_w = Inches(3.7)
for i, (timing, title, body) in enumerate(stappen):
    sx = Inches(0.5) + i * Inches(4.25)
    sy = Inches(1.75)

    # Hoofdkaart
    card(sl, sx, sy, step_w, Inches(5.0),
         fill=RGBColor(0xf5, 0xf7, 0xff),
         border=RGBColor(0xc7, 0xd4, 0xfd))

    # Nummer cirkel bovenaan
    circ(sl, sx + Inches(1.55), sy + Inches(0.3), Inches(0.6), BLUE)
    txt(sl, str(i + 1),
        sx + Inches(1.55), sy + Inches(0.3), Inches(0.6), Inches(0.6),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Timing badge
    txt(sl, timing,
        sx + Inches(0.2), sy + Inches(1.1), step_w - Inches(0.4), Inches(0.55),
        size=11, color=MUTED, align=PP_ALIGN.CENTER)

    txt(sl, title,
        sx + Inches(0.2), sy + Inches(1.7), step_w - Inches(0.4), Inches(0.55),
        size=17, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    txt(sl, body,
        sx + Inches(0.25), sy + Inches(2.4), step_w - Inches(0.5), Inches(2.3),
        size=13, color=MUTED)

    # Pijl rechts
    if i < 2:
        txt(sl, "→",
            sx + step_w + Inches(0.1), sy + Inches(2.2), Inches(0.35), Inches(0.6),
            size=24, bold=True, color=BLUE)

txt(sl, "Geen abonnementsrisico: de eerste maand is gratis. Daarna maandelijks opzegbaar.",
    Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.38),
    size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — INVESTERING
# Boodschap: Eenvoudig. Transparant. Maandelijks opzegbaar.
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Investering", "Transparant — geen setup-kost, geen lange contracten")

# Aanbevolen plan voor Waterfront (groot, centraal)
rect(sl, Inches(3.5), Inches(1.65), Inches(6.3), Inches(5.3), fill=BLUE)

# Badge
rect(sl, Inches(5.1), Inches(1.38), Inches(3.1), Inches(0.42), fill=GREEN)
txt(sl, "Aanbevolen voor Waterfront Genk",
    Inches(5.1), Inches(1.42), Inches(3.1), Inches(0.36),
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Groei",
    Inches(3.75), Inches(1.85), Inches(5.8), Inches(0.6),
    size=24, bold=True, color=WHITE)
txt(sl, "€ 249",
    Inches(3.75), Inches(2.5), Inches(3.5), Inches(1.0),
    size=52, bold=True, color=WHITE)
txt(sl, "/ maand  (ex. btw)",
    Inches(6.0), Inches(3.1), Inches(3.5), Inches(0.45),
    size=15, color=RGBColor(0xc7, 0xd4, 0xfd))

inclusions = [
    "1 locatie — Waterfront Genk",
    "14-daagse dagelijkse forecast",
    "Personeelsplanning inbegrepen",
    "Onboarding & databegeleiding",
    "Priority support",
    "Maandelijks opzegbaar",
]
for i, inc in enumerate(inclusions):
    iy = Inches(3.7) + i * Inches(0.48)
    dot(sl, Inches(3.8), iy + Inches(0.1), color=WHITE)
    txt(sl, inc, Inches(4.0), iy, Inches(5.5), Inches(0.42),
        size=13, color=WHITE)

# Kleine kaartjes links en rechts
for side, (plan_name, price, items) in enumerate([
    ("Starter", "€ 149/maand", ["1 locatie", "14-daagse forecast", "Upload via CSV", "E-mail support"]),
    ("Enterprise", "Op maat", ["Meerdere locaties", "API-koppeling", "Custom rapportage", "SLA-garantie"]),
]):
    sx = Inches(0.4) if side == 0 else Inches(10.05)
    card(sl, sx, Inches(2.0), Inches(2.8), Inches(4.6),
         fill=RGBColor(0xf8, 0xf9, 0xff))
    txt(sl, plan_name, sx + Inches(0.2), Inches(2.15), Inches(2.4), Inches(0.4),
        size=14, bold=True, color=TEXT)
    txt(sl, price, sx + Inches(0.2), Inches(2.6), Inches(2.4), Inches(0.55),
        size=18, bold=True, color=MUTED)
    for j, item in enumerate(items):
        dot(sl, sx + Inches(0.22), Inches(3.35) + j * Inches(0.5), color=MUTED)
        txt(sl, item, sx + Inches(0.42), Inches(3.28) + j * Inches(0.5),
            Inches(2.2), Inches(0.42), size=12, color=MUTED)

txt(sl, "Eerste maand gratis. Daarna maandelijks. Jaarlijkse betaling: 2 maanden gratis.",
    Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.38),
    size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CTA / AFSLUITING
# Boodschap: Één vraag. Starten we volgende week?
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

# Donkerblauwe achtergrond
rect(sl, 0, 0, W, H, fill=RGBColor(0x0f, 0x2a, 0x9e))

# Subtiele decoratieve vlakken
rect(sl, Inches(8.5), 0, Inches(5), H, fill=RGBColor(0x1a, 0x44, 0xe8))
rect(sl, Inches(10), Inches(1.5), Inches(4), Inches(4.5),
     fill=RGBColor(0x2d, 0x55, 0xf0))

txt(sl, "Kris,",
    Inches(0.7), Inches(1.0), Inches(8), Inches(0.75),
    size=40, bold=True, color=WHITE)

txt(sl, "starten we volgende week?",
    Inches(0.7), Inches(1.8), Inches(8.5), Inches(1.1),
    size=38, color=WHITE)

rect(sl, Inches(0.7), Inches(3.1), Inches(4), Pt(3),
     fill=RGBColor(0xc7, 0xd4, 0xfd))

txt(sl, "We uploaden je historische data, het model leert een week,\n"
        "en daarna zie je elke ochtend wat de dag gaat brengen.\n"
        "Eerste maand gratis — geen risico.",
    Inches(0.7), Inches(3.25), Inches(8.0), Inches(1.5),
    size=16, color=RGBColor(0xc7, 0xd4, 0xfd))

# Drie snelle punten
for i, punt in enumerate([
    "Eerste maand gratis",
    "Klaar in 1 week",
    "Maandelijks opzegbaar",
]):
    circ(sl, Inches(0.7) + i * Inches(3.0), Inches(5.0), Inches(0.28),
         fill=GREEN)
    txt(sl, punt,
        Inches(1.1) + i * Inches(3.0), Inches(4.95), Inches(2.7), Inches(0.38),
        size=13, bold=True, color=WHITE)

# Contact onderaan
txt(sl, "Sven Schaekers  |  CloudCast Analytics",
    Inches(0.7), Inches(6.2), Inches(8), Inches(0.4),
    size=14, bold=True, color=WHITE)
txt(sl, "schaekers.sven@gmail.com  |  cloudcastanalytics.com",
    Inches(0.7), Inches(6.6), Inches(8), Inches(0.38),
    size=13, color=RGBColor(0xc7, 0xd4, 0xfd))


# ── Opslaan ────────────────────────────────────────────────────────────────
output = r"C:\Users\eddys\OneDrive (1)\Zakelijke documenten\CloudCast_Waterfront_Genk.pptx"
prs.save(output)
print(f"Klaar: {output}")
