"""
CloudCast Analytics — Presentatie generator
Genereert een professionele .pptx met CloudCast branding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Kleuren ────────────────────────────────────────────────────────────────
BLUE        = RGBColor(0x1a, 0x44, 0xe8)
BLUE_DARK   = RGBColor(0x0f, 0x2a, 0x9e)
BLUE_MID    = RGBColor(0x2d, 0x55, 0xf0)
BLUE_LIGHT  = RGBColor(0xe8, 0xee, 0xff)
TEXT_DARK   = RGBColor(0x1a, 0x1f, 0x36)
TEXT_MUTED  = RGBColor(0x6b, 0x72, 0x80)
WHITE       = RGBColor(0xff, 0xff, 0xff)
BG_LIGHT    = RGBColor(0xf5, 0xf7, 0xff)
GREEN       = RGBColor(0x16, 0xa3, 0x4a)
AMBER       = RGBColor(0xd9, 0x77, 0x06)

# ── Afmetingen (16:9 widescreen) ───────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # volledig leeg


# ── Hulpfuncties ───────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, font_name="Calibri", italic=False,
             word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, x, y, w, h,
                  font_size=16, color=TEXT_DARK, line_spacing=1.15,
                  font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def slide_header(slide, title, subtitle=None):
    """Blauwe balk bovenaan met witte titel."""
    bar = add_rect(slide, 0, 0, W, Inches(1.4), fill_color=BLUE)
    add_text(slide, title,
             Inches(0.6), Inches(0.22), Inches(11), Inches(0.65),
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(0.85), Inches(11), Inches(0.4),
                 font_size=14, color=RGBColor(0xc7, 0xd4, 0xfd))


def kpi_card(slide, x, y, w, h, value, label, sublabel=None,
             value_color=BLUE):
    """Klein KPI-kaartje met witte achtergrond en rand."""
    card = add_rect(slide, x, y, w, h, fill_color=WHITE)
    card.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
    card.line.width = Pt(1)
    add_text(slide, value,
             x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.55),
             font_size=32, bold=True, color=value_color)
    add_text(slide, label,
             x + Inches(0.2), y + Inches(0.65), w - Inches(0.4), Inches(0.3),
             font_size=13, bold=True, color=TEXT_DARK)
    if sublabel:
        add_text(slide, sublabel,
                 x + Inches(0.2), y + Inches(0.9), w - Inches(0.4), Inches(0.3),
                 font_size=11, color=TEXT_MUTED)


def bullet_block(slide, x, y, w, items, title=None,
                 font_size=15, dot_color=BLUE):
    """Lijst met blauwe bullets."""
    cur_y = y
    if title:
        add_text(slide, title, x, cur_y, w, Inches(0.35),
                 font_size=14, bold=True, color=TEXT_DARK)
        cur_y += Inches(0.38)
    for item in items:
        # blauwe bol
        dot = slide.shapes.add_shape(9, x, cur_y + Inches(0.06),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        add_text(slide, item, x + Inches(0.22), cur_y,
                 w - Inches(0.22), Inches(0.32),
                 font_size=font_size, color=TEXT_DARK)
        cur_y += Inches(0.36)
    return cur_y


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Achtergrond: lichtblauw gradient nabootsing (één rechthoek)
bg = add_rect(slide, 0, 0, W, H, fill_color=BG_LIGHT)

# Blauwe accent-balk links
accent = add_rect(slide, 0, 0, Inches(0.45), H, fill_color=BLUE)

# Blok rechts met illustratieve donkere rechthoek (dashboard mockup sfeer)
dark_block = add_rect(slide, Inches(7.8), Inches(1.2), Inches(5.1), Inches(5.1),
                      fill_color=RGBColor(0xf0, 0xf4, 0xff))
dark_block.line.color.rgb = RGBColor(0xc7, 0xd4, 0xfd)
dark_block.line.width = Pt(1.5)

# Nep-dashboard blokjes binnen het schermvak
for i, (row_y, col_x, w_b, h_b, col) in enumerate([
    (Inches(1.5),  Inches(8.0),  Inches(2.2), Inches(0.9),  RGBColor(0xe8, 0xee, 0xff)),
    (Inches(1.5),  Inches(10.3), Inches(2.4), Inches(0.9),  RGBColor(0xe8, 0xee, 0xff)),
    (Inches(2.55), Inches(8.0),  Inches(4.7), Inches(2.1),  WHITE),
    (Inches(4.8),  Inches(8.0),  Inches(2.2), Inches(1.2),  WHITE),
    (Inches(4.8),  Inches(10.3), Inches(2.4), Inches(1.2),  WHITE),
]):
    b = add_rect(slide, col_x, row_y, w_b, h_b, fill_color=col)
    b.line.color.rgb = RGBColor(0xc7, 0xd4, 0xfd)
    b.line.width = Pt(0.75)

# Grafiek-nep-balken
bar_data = [0.4, 0.6, 0.5, 0.8, 0.7, 0.95, 1.0]
bar_w = Inches(0.42)
for i, pct in enumerate(bar_data):
    bar_h = Inches(1.5 * pct)
    bx = Inches(8.1) + i * Inches(0.52)
    by = Inches(2.55) + Inches(2.1) - bar_h - Inches(0.05)
    b = add_rect(slide, bx, by, bar_w, bar_h, fill_color=BLUE_MID)
    b.line.fill.background()

# Logo-tekst (links, verticaal gecentreerd)
add_text(slide, "CloudCast",
         Inches(0.65), Inches(1.5), Inches(5.5), Inches(1.0),
         font_size=46, bold=True, color=TEXT_DARK, font_name="Calibri")
add_text(slide, "Analytics",
         Inches(0.65), Inches(2.4), Inches(5.5), Inches(0.8),
         font_size=46, bold=False, color=BLUE, font_name="Calibri")

# Divider lijn
div = add_rect(slide, Inches(0.65), Inches(3.3), Inches(3.5), Pt(3),
               fill_color=BLUE)

add_text(slide, "Voorspel je drukte. Wees voorbereid.",
         Inches(0.65), Inches(3.5), Inches(6.5), Inches(0.5),
         font_size=17, color=TEXT_MUTED, italic=True)

add_text(slide, "Revenue Forecasting & Personeelsplanning voor Horeca & Leisure",
         Inches(0.65), Inches(4.1), Inches(6.8), Inches(0.55),
         font_size=13, color=TEXT_MUTED)

# Datum
add_text(slide, "Juni 2026",
         Inches(0.65), Inches(6.6), Inches(3), Inches(0.4),
         font_size=12, color=TEXT_MUTED)

add_text(slide, "cloudcastanalytics.com",
         Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.4),
         font_size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — HET PROBLEEM
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill_color=WHITE)
slide_header(slide, "Het probleem", "Waarom planning in horeca & leisure structureel mislukt")

# Drie pijnpunten
pain_items = [
    ("Planning op gevoel",
     "Hoe druk wordt het morgen? Volgende zaterdag? Niemand weet het echt. "
     "Roosters worden opgesteld op basis van buikgevoel of vorig jaar."),
    ("Overbezetting kost geld",
     "Te veel personeel inplannen op een rustige dag betekent directe loonkost "
     "zonder omzet. Bij een gemiddelde horeca-zaak loopt dit op tot €15.000+ per jaar."),
    ("Onderbezetting kost klanten",
     "Te weinig personeel op een drukke dag geeft lange wachttijden, "
     "gefrustreerde klanten en verloren omzet. Eén slechte ervaring = minder reviews."),
]

card_w = Inches(3.6)
card_h = Inches(4.5)
for i, (title, body) in enumerate(pain_items):
    cx = Inches(0.5) + i * Inches(4.22)
    cy = Inches(1.7)
    card = add_rect(slide, cx, cy, card_w, card_h,
                    fill_color=RGBColor(0xf8, 0xf9, 0xff))
    card.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
    card.line.width = Pt(1)
    # Nummer
    num_bg = add_rect(slide, cx + Inches(0.25), cy + Inches(0.25),
                      Inches(0.5), Inches(0.5), fill_color=BLUE)
    add_text(slide, str(i + 1),
             cx + Inches(0.25), cy + Inches(0.25),
             Inches(0.5), Inches(0.5),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title,
             cx + Inches(0.25), cy + Inches(0.9), card_w - Inches(0.5), Inches(0.55),
             font_size=16, bold=True, color=TEXT_DARK)
    add_text(slide, body,
             cx + Inches(0.25), cy + Inches(1.55), card_w - Inches(0.5), Inches(2.6),
             font_size=13, color=TEXT_MUTED)

# Stat onderaan
stat_bg = add_rect(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7),
                   fill_color=BLUE_LIGHT)
stat_bg.line.color.rgb = RGBColor(0xc7, 0xd4, 0xfd)
stat_bg.line.width = Pt(1)
add_text(slide,
         "60–70% nauwkeurigheid bij plannen op gevoel  |  "
         "CloudCast haalt gemiddeld 97%  |  "
         "Dat verschil vertaalt zich direct in minder personeelskost",
         Inches(0.8), Inches(6.48), Inches(12), Inches(0.5),
         font_size=12, color=BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DE OPLOSSING
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill_color=WHITE)
slide_header(slide, "De oplossing", "CloudCast Analytics — van historische data naar dagelijkse forecast")

# Pijlen flow: DATA → MODEL → FORECAST → ACTIE
steps = [
    ("1", "Historische\ndata uploaden", "CSV of directe\nkoppeling kassasysteem"),
    ("2", "Model leert\npatronen", "Seizoenen, weekdagen,\nweer & evenementen"),
    ("3", "Dagelijkse\nforecast", "Omzet & bezoekers\nvoor 14 dagen vooruit"),
    ("4", "Personeels-\nadvies", "Exact hoeveel FTE\nper dag nodig"),
]

step_w = Inches(2.8)
for i, (num, title, body) in enumerate(steps):
    cx = Inches(0.4) + i * Inches(3.2)
    cy = Inches(1.8)

    # Kaart
    card = add_rect(slide, cx, cy, step_w, Inches(3.8),
                    fill_color=RGBColor(0xf5, 0xf7, 0xff))
    card.line.color.rgb = RGBColor(0xc7, 0xd4, 0xfd)
    card.line.width = Pt(1.5)

    # Stap cirkel
    circ = slide.shapes.add_shape(9, cx + Inches(1.1), cy + Inches(0.25),
                                   Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    add_text(slide, num,
             cx + Inches(1.1), cy + Inches(0.25),
             Inches(0.6), Inches(0.6),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, title,
             cx + Inches(0.2), cy + Inches(1.05), step_w - Inches(0.4), Inches(0.8),
             font_size=15, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, body,
             cx + Inches(0.2), cy + Inches(1.9), step_w - Inches(0.4), Inches(1.6),
             font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Pijl naar rechts (behalve laatste)
    if i < 3:
        arr = add_rect(slide, cx + step_w + Inches(0.05), cy + Inches(1.7),
                       Inches(0.3), Pt(3), fill_color=BLUE)

# Ondertekst
add_text(slide,
         "Geen technische kennis vereist. Klaar in 1 dag. Werkt met elk kassasysteem of Excel-export.",
         Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
         font_size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WAT JE KRIJGT (FEATURES)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill_color=WHITE)
slide_header(slide, "Wat je krijgt", "Alles in één platform — zonder technische kennis")

features = [
    ("14-daagse forecast",
     ["Dagelijkse voorspelling van omzet en bezoekers",
      "Weerseffect en seizoenspatronen ingebouwd",
      "Betrouwbaarheidsinterval per dag"]),
    ("Personeelsplanning",
     ["Aanbevolen FTE per dag op basis van forecast",
      "Instelbare personeelsregels per locatie",
      "Directe koppeling met drukteniveau"]),
    ("Data management",
     ["CSV of Excel upload in minuten",
      "Automatische kolomherkenning",
      "Validatie en foutrapportage"]),
    ("Multi-locatie",
     ["Meerdere vestigingen in één platform",
      "Per locatie eigen regels en forecast",
      "Gecentraliseerd overzicht"]),
]

feat_w = Inches(5.9)
feat_h = Inches(2.35)
for i, (title, bullets) in enumerate(features):
    col = i % 2
    row = i // 2
    cx = Inches(0.45) + col * Inches(6.4)
    cy = Inches(1.65) + row * Inches(2.6)

    card = add_rect(slide, cx, cy, feat_w, feat_h,
                    fill_color=RGBColor(0xf8, 0xf9, 0xff))
    card.line.color.rgb = RGBColor(0xc7, 0xd4, 0xfd)
    card.line.width = Pt(1)

    # Blauwe bovenbalk
    top_bar = add_rect(slide, cx, cy, feat_w, Inches(0.38), fill_color=BLUE)

    add_text(slide, title,
             cx + Inches(0.2), cy + Inches(0.05), feat_w - Inches(0.4), Inches(0.3),
             font_size=13, bold=True, color=WHITE)

    for j, bullet in enumerate(bullets):
        dot = slide.shapes.add_shape(9,
                                     cx + Inches(0.2),
                                     cy + Inches(0.55) + j * Inches(0.52),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_text(slide, bullet,
                 cx + Inches(0.38),
                 cy + Inches(0.48) + j * Inches(0.52),
                 feat_w - Inches(0.55), Inches(0.45),
                 font_size=12, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROI & BESPARING
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill_color=WHITE)
slide_header(slide, "Wat levert het op?", "Concrete besparing op jaarbasis")

# Grote blauwe kaart links
big_card = add_rect(slide, Inches(0.45), Inches(1.7), Inches(5.5), Inches(4.5),
                    fill_color=BLUE)
add_text(slide, "Geschatte jaarlijkse besparing",
         Inches(0.7), Inches(2.0), Inches(5.0), Inches(0.5),
         font_size=14, color=RGBColor(0xc7, 0xd4, 0xfd))
add_text(slide, "€ 17.000 – 28.000",
         Inches(0.7), Inches(2.55), Inches(5.0), Inches(1.1),
         font_size=36, bold=True, color=WHITE)
add_text(slide,
         "Op basis van gemiddeld 2 FTE minder overbezetting\n"
         "op 60% van de weekenddagen en 40% van de weekdagen.\n"
         "Loonkost €110–€150 per medewerker per dag.",
         Inches(0.7), Inches(3.75), Inches(5.0), Inches(1.8),
         font_size=12, color=RGBColor(0xc7, 0xd4, 0xfd))
add_text(slide, "Terugverdientijd: 1 tot 3 maanden",
         Inches(0.7), Inches(5.6), Inches(5.0), Inches(0.4),
         font_size=13, bold=True, color=RGBColor(0x86, 0xef, 0xac))

# Drie KPI-kaarten rechts
kpi_data = [
    ("97%",       "Voorspellingsnauwkeurigheid",  "vs. 60-70% op gevoel",   BLUE),
    ("-3 uur",    "Minder planningswerk/week",     "Geen handmatige Excel",  GREEN),
    ("14 dagen",  "Forecast horizon",              "Per dag, per locatie",   RGBColor(0x71, 0x35, 0xd4)),
]

for i, (val, lbl, sub, col) in enumerate(kpi_data):
    kx = Inches(6.4)
    ky = Inches(1.7) + i * Inches(1.65)
    kpi_card(slide, kx, ky, Inches(6.45), Inches(1.45),
             val, lbl, sub, value_color=col)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — VOOR WIE?
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill_color=WHITE)
slide_header(slide, "Voor wie is CloudCast?", "Ontworpen voor horeca, leisure en evenementenlocaties")

sectors = [
    ("Horeca",        ["Restaurants", "Cafés & brasserieën", "Foodtrucks & pop-ups"]),
    ("Leisure",       ["Zwem- en waterparken", "Sportcentra", "Pretparken"]),
    ("Events",        ["Evenementenzalen", "Festivals", "Markten & beurzen"]),
    ("Retail",        ["Winkels met seizoenspieken", "Pop-up stores", "Flagship stores"]),
]

sec_w = Inches(2.8)
for i, (sector, items) in enumerate(sectors):
    cx = Inches(0.45) + i * Inches(3.2)
    cy = Inches(1.8)

    card = add_rect(slide, cx, cy, sec_w, Inches(4.2),
                    fill_color=RGBColor(0xf5, 0xf7, 0xff))
    card.line.color.rgb = RGBColor(0xc7, 0xd4, 0xfd)
    card.line.width = Pt(1)

    top = add_rect(slide, cx, cy, sec_w, Inches(0.55), fill_color=BLUE_LIGHT)
    add_text(slide, sector,
             cx + Inches(0.2), cy + Inches(0.1), sec_w - Inches(0.4), Inches(0.38),
             font_size=15, bold=True, color=BLUE)

    for j, item in enumerate(items):
        dot = slide.shapes.add_shape(9,
                                     cx + Inches(0.2),
                                     cy + Inches(0.8) + j * Inches(0.65),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_text(slide, item,
                 cx + Inches(0.38),
                 cy + Inches(0.72) + j * Inches(0.65),
                 sec_w - Inches(0.55), Inches(0.55),
                 font_size=13, color=TEXT_DARK)

add_text(slide,
         "Elke sector waar drukte fluctueert en personeel de grootste kostenpost is.",
         Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.6),
         font_size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ONBOARDING / STAPPENPLAN
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill_color=WHITE)
slide_header(slide, "Hoe kom je aan de slag?", "Van handdruk tot eerste forecast in 3 stappen")

stappen = [
    ("Stap 1 — Week 1",    "Koppeling & data",
     "We verbinden CloudCast met jouw kassasysteem of je uploadt een Excel-export. "
     "Eenmalige instelling — wij begeleiden je er volledig doorheen."),
    ("Stap 2 — Week 2–3",  "Leerperiode",
     "Het model leert van jouw historische data. Na 2 weken zijn de eerste "
     "voorspellingen al bruikbaar. Na 4 weken zit je op 90%+ nauwkeurigheid."),
    ("Stap 3 — Dag 1",     "Dagelijks gebruik",
     "Elke ochtend een nieuwe forecast voor de komende 14 dagen. "
     "Personeelsadvies per dag. Volledig automatisch — geen actie van jou vereist."),
]

for i, (timing, title, body) in enumerate(stappen):
    cy = Inches(1.7) + i * Inches(1.7)

    # Tijdslijn-lijn
    line = add_rect(slide, Inches(1.45), cy + Inches(0.65),
                    Inches(10.4), Pt(2), fill_color=RGBColor(0xe2, 0xe8, 0xf0))

    # Cirkel
    circ = slide.shapes.add_shape(9, Inches(0.45), cy + Inches(0.35),
                                   Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    add_text(slide, str(i + 1),
             Inches(0.45), cy + Inches(0.35),
             Inches(0.6), Inches(0.6),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, timing,
             Inches(1.3), cy + Inches(0.05), Inches(2.5), Inches(0.35),
             font_size=11, color=TEXT_MUTED, bold=False)
    add_text(slide, title,
             Inches(1.3), cy + Inches(0.38), Inches(4), Inches(0.45),
             font_size=15, bold=True, color=TEXT_DARK)
    add_text(slide, body,
             Inches(1.3), cy + Inches(0.85), Inches(10.5), Inches(0.65),
             font_size=13, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRIJZEN (INVESTERING)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill_color=WHITE)
slide_header(slide, "Investering", "Transparante prijsstructuur — geen verrassingen")

plannen = [
    ("Starter",   "€ 149 / maand",  False, [
        "1 locatie",
        "14-daagse forecast",
        "CSV / Excel upload",
        "Personeelsplanning",
        "E-mail support",
    ]),
    ("Groei",     "€ 249 / maand",  True,  [
        "Tot 3 locaties",
        "14-daagse forecast",
        "Kassasysteem koppeling",
        "Personeelsplanning",
        "Maandrapport automatisch",
        "Priority support",
    ]),
    ("Enterprise", "Op maat",        False, [
        "Onbeperkt locaties",
        "API-integraties",
        "Dedicated onboarding",
        "SLA-garantie",
        "Custom rapportage",
    ]),
]

plan_w = Inches(3.6)
for i, (name, price, highlighted, items) in enumerate(plannen):
    cx = Inches(0.5) + i * Inches(4.2)
    cy = Inches(1.65)
    ph = Inches(5.1)

    if highlighted:
        card = add_rect(slide, cx, cy, plan_w, ph, fill_color=BLUE)
        txt_color = WHITE
        sub_color = RGBColor(0xc7, 0xd4, 0xfd)
        dot_color = WHITE
        border_col = BLUE
    else:
        card = add_rect(slide, cx, cy, plan_w, ph,
                        fill_color=RGBColor(0xf8, 0xf9, 0xff))
        card.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
        card.line.width = Pt(1)
        txt_color = TEXT_DARK
        sub_color = TEXT_MUTED
        dot_color = BLUE
        border_col = None

    if highlighted:
        badge = add_rect(slide, cx + Inches(0.7), cy - Inches(0.18),
                         Inches(2.2), Inches(0.36),
                         fill_color=GREEN)
        add_text(slide, "Meest gekozen",
                 cx + Inches(0.7), cy - Inches(0.18),
                 Inches(2.2), Inches(0.36),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, name,
             cx + Inches(0.25), cy + Inches(0.2), plan_w - Inches(0.5), Inches(0.4),
             font_size=15, bold=True, color=txt_color)
    add_text(slide, price,
             cx + Inches(0.25), cy + Inches(0.65), plan_w - Inches(0.5), Inches(0.65),
             font_size=26, bold=True, color=txt_color)

    for j, item in enumerate(items):
        dot = slide.shapes.add_shape(9,
                                     cx + Inches(0.25),
                                     cy + Inches(1.5) + j * Inches(0.6),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        add_text(slide, item,
                 cx + Inches(0.45),
                 cy + Inches(1.42) + j * Inches(0.6),
                 plan_w - Inches(0.65), Inches(0.52),
                 font_size=12, color=txt_color)

add_text(slide,
         "Alle plannen inclusief onboarding. Jaarlijkse betaling geeft 2 maanden gratis.",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
         font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — VOLGENDE STAP (CTA)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Volledig blauwe achtergrond
add_rect(slide, 0, 0, W, H, fill_color=BLUE)

# Subtiele decoratieve rechthoeken
add_rect(slide, Inches(9.5), 0, Inches(4), Inches(7.5),
         fill_color=RGBColor(0x2d, 0x55, 0xf0))
add_rect(slide, Inches(11), Inches(2), Inches(3), Inches(3),
         fill_color=RGBColor(0x1a, 0x3a, 0xd0))

add_text(slide, "CloudCast Analytics",
         Inches(0.8), Inches(1.2), Inches(9), Inches(0.6),
         font_size=16, color=RGBColor(0xc7, 0xd4, 0xfd))

add_text(slide, "Klaar om slimmer\nte plannen?",
         Inches(0.8), Inches(1.9), Inches(9), Inches(2.0),
         font_size=48, bold=True, color=WHITE, font_name="Calibri")

add_text(slide, "Plan een gratis demo van 30 minuten. We laten je zien hoe CloudCast werkt\n"
                "met jouw eigen data — geen verplichtingen, geen technische kennis vereist.",
         Inches(0.8), Inches(4.1), Inches(8.5), Inches(1.2),
         font_size=16, color=RGBColor(0xc7, 0xd4, 0xfd))

# CTA knop nep
cta = add_rect(slide, Inches(0.8), Inches(5.5), Inches(3.5), Inches(0.75), fill_color=WHITE)
add_text(slide, "Plan een gratis demo",
         Inches(0.8), Inches(5.55), Inches(3.5), Inches(0.65),
         font_size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Contact
contact_items = [
    "cloudcastanalytics.com",
    "info@cloudcastanalytics.com",
]
for i, item in enumerate(contact_items):
    add_text(slide, item,
             Inches(0.8), Inches(6.4) + i * Inches(0.35),
             Inches(6), Inches(0.32),
             font_size=13, color=RGBColor(0xc7, 0xd4, 0xfd))


# ── Opslaan ────────────────────────────────────────────────────────────────
output_path = r"C:\Users\eddys\OneDrive (1)\Zakelijke documenten\CloudCast_Presentatie.pptx"
prs.save(output_path)
print(f"Presentatie opgeslagen: {output_path}")
